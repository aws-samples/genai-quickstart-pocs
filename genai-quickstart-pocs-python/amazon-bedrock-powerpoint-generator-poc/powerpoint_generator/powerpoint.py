import json
from uuid import uuid4
import logging
import sys
from pptx import Presentation
import os
from .models import CompleteSectionWithSlides

# This is specific to the template.pptx provided in solution.
# Changing the template or the MasterSlides in the template will likely require these fields to update.
TITLE_SLIDE_LAYOUT = 1
CONTENT_SLIDE_LAYOUT = 0
SECTION_HEADER_SLIDE_LAYOUT = 2
AGENDA_SLIDE_LAYOUT = 0

logger = logging.getLogger(__name__)


def generate_powerpoint_file(topic: str, presentation_content: list[CompleteSectionWithSlides], write_callback=None) -> str:
    """
    Generates the PowerPoint Presentation file using the generated and researched content.
    """
    logger.info("Generating PowerPoint file")
    if write_callback:
        write_callback("Generating PowerPoint file")
    preso = Presentation(os.path.join(os.path.dirname(__file__), "template.pptx"))
    content_layout = preso.slide_layouts[CONTENT_SLIDE_LAYOUT]
    title_slide_layout = preso.slide_layouts[TITLE_SLIDE_LAYOUT]
    agenda_slide_layout = preso.slide_layouts[AGENDA_SLIDE_LAYOUT]
    section_header_slide_layout = preso.slide_layouts[SECTION_HEADER_SLIDE_LAYOUT]
    logger.info("Adding title slide")
    title_slide = preso.slides.add_slide(title_slide_layout)
    title_slide.shapes.placeholders[0].text = topic
    logger.info("Adding agenda slide")
    agenda_items = [section.title for section in presentation_content]
    agenda_slide = preso.slides.add_slide(agenda_slide_layout)
    agenda_slide.shapes.title.text = "Agenda"
    agenda_slide.shapes.placeholders[1].text = "\n".join(agenda_items)
    for section in presentation_content:
        logger.info(f"Adding section headers: {section.title}")
        section_header_slide = preso.slides.add_slide(section_header_slide_layout)
        section_header_slide.shapes.title.text = section.title
        for slide in section.slides:
            logger.info(f"Adding slide: {slide.title}")
            content_slide = preso.slides.add_slide(content_layout)
            content_slide.shapes.title.text = topic
            content_slide.shapes.placeholders[0].text = slide.title
            content_slide.shapes.placeholders[1].text = "\n".join(slide.main_content.split(">>"))
            notes_paragraph = content_slide.notes_slide.notes_text_frame.add_paragraph()
            notes_paragraph.text = slide.presenter_notes

    conclusion_slide = preso.slides.add_slide(section_header_slide_layout)
    conclusion_slide.shapes.placeholders[0].text = "Thank You!"
    file_path = os.path.join(os.path.dirname(__file__), f"../temp/{uuid4()}.pptx")
    preso.save(file_path)
    logger.info("PowerPoint file generated successfully!")
    if write_callback:
        write_callback("PowerPoint file generated successfully!")
    logger.info(f"File saved to: {file_path}")
    return file_path

def delete_file(file_path: str):
    try:
        os.remove(file_path)
    except Exception as e:
        logger.error(f"Error deleting file: {e}")